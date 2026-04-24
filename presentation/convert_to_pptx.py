#!/usr/bin/env python3
"""
convert_to_pptx.py
Build presentation_anchoveta.pptx from the Beamer source.
Run from: /home/jupyter-daniela/peru_catch_modeling/
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as PILImage

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE    = Path(__file__).resolve().parent
IMGDIR  = BASE / "outputs"
OUT     = BASE / "presentation_anchoveta.pptx"
FOOTER  = "Suyana  |  Seguro Paramétrico Anchoveta Perú  |  Abril 2026"

# ── Slide dimensions ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Colors ────────────────────────────────────────────────────────────────────
C_GREEN   = RGBColor(0x43, 0xA0, 0x47)
C_BLACK   = RGBColor(0x14, 0x14, 0x14)
C_DARK    = RGBColor(0x11, 0x11, 0x11)
C_GRAY    = RGBColor(0x55, 0x55, 0x55)
C_MIDGRAY = RGBColor(0x9E, 0x9E, 0x9E)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Zone constants ────────────────────────────────────────────────────────────
RULE_Y       = Inches(0)
RULE_H       = Emu(38100)
TITLE_Y      = Inches(0.38)
TITLE_H      = Inches(0.40)
SUBT_Y       = Inches(0.78)
SUBT_H       = Inches(0.28)
TRULE_Y_NO_SUBT = Inches(0.78)
TRULE_Y_SUBT    = Inches(1.06)
CONTENT_Y_NO_SUBT = Inches(0.88)
CONTENT_Y_SUBT    = Inches(1.16)
CONTENT_H_NO_SUBT = Inches(6.12)
CONTENT_H_SUBT    = Inches(5.84)
CONTENT_L    = Inches(0.45)
CONTENT_W    = Inches(12.43)
FOOTER_RULE_Y = Inches(7.10)
FOOTER_Y      = Inches(7.14)
FOOTER_H      = Inches(0.36)

TOTAL = 27


# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rule(slide, y, color, height_emu=Emu(22860), width=None):
    w = width or SLIDE_W
    shape = slide.shapes.add_shape(1, 0, y, w, height_emu)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size_pt, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or C_BLACK
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = align
    return txb


def add_multiline(slide, lines, x, y, w, h, size_pt, color=None, bold=False,
                  italic=False, leading_pt=None):
    """lines: list of (text, bold, color, italic) or plain str."""
    color = color or C_BLACK
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, lb, lc, li = line, bold, color, italic
        else:
            txt, lb, lc, li = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if leading_pt:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPts = etree.SubElement(lnSpc, qn("a:spcPts"))
            spcPts.set("val", str(int(leading_pt * 100)))
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size_pt)
        run.font.bold = lb
        run.font.italic = li
        run.font.color.rgb = lc if lc else color
    return txb


def add_header(slide):
    add_rule(slide, RULE_Y, C_GREEN, height_emu=RULE_H)
    # Logo placeholder (no brand files present)
    add_text(slide, "SUYANA", SLIDE_W - Inches(0.90), Inches(0.02),
             Inches(0.80), Inches(0.22), size_pt=7, bold=True,
             color=C_GREEN, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_num):
    add_rule(slide, FOOTER_RULE_Y, C_MIDGRAY, height_emu=Emu(12700))
    add_text(slide, FOOTER, Inches(0.18), FOOTER_Y, Inches(10), FOOTER_H,
             size_pt=7, color=C_MIDGRAY)
    add_text(slide, f"{page_num} / {TOTAL}", Inches(0.18), FOOTER_Y,
             SLIDE_W - Inches(0.36), FOOTER_H, size_pt=7, color=C_MIDGRAY,
             align=PP_ALIGN.RIGHT)


def add_frametitle(slide, title, subtitle=None):
    add_text(slide, title, CONTENT_L, TITLE_Y, CONTENT_W, TITLE_H,
             size_pt=18, bold=True, color=C_BLACK)
    if subtitle:
        add_text(slide, subtitle, CONTENT_L, SUBT_Y, CONTENT_W, SUBT_H,
                 size_pt=11, color=C_GRAY)
        rule_y   = TRULE_Y_SUBT
        content_y = CONTENT_Y_SUBT
    else:
        rule_y   = TRULE_Y_NO_SUBT
        content_y = CONTENT_Y_NO_SUBT
    add_rule(slide, rule_y, C_MIDGRAY, height_emu=Emu(12700))
    return content_y


def img_dims(path, max_w_emu, max_h_emu):
    with PILImage.open(path) as im:
        pw, ph = im.size
    ratio = min(max_w_emu / pw, max_h_emu / ph)
    return int(pw * ratio), int(ph * ratio)


def add_image_centered(slide, img_path, content_top, cap_h_emu=0, max_frac=1.0):
    """Center image below content_top, reserving cap_h_emu at the bottom."""
    avail_h = FOOTER_RULE_Y - content_top - cap_h_emu - Inches(0.10)
    avail_w = CONTENT_W * max_frac
    iw, ih  = img_dims(img_path, avail_w, avail_h)
    left    = CONTENT_L + (CONTENT_W - iw) // 2
    top     = content_top + Emu(40000)
    slide.shapes.add_picture(str(img_path), left, top, iw, ih)
    return top, iw, ih


def add_caption(slide, text, top, ih, size_pt=8):
    cap_y = top + ih + Emu(30000)
    add_text(slide, text, CONTENT_L, cap_y, CONTENT_W, Inches(0.55),
             size_pt=size_pt, color=C_BLACK, wrap=True)


# ── Dark slide builders ───────────────────────────────────────────────────────
def title_slide(prs, page_num):
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_text(slide, "SUYANA", Inches(0.45), Inches(1.00), Inches(3), Inches(0.40),
             size_pt=14, bold=True, color=C_GREEN)
    add_text(slide, "Seguro Paramétrico de Captura de Anchoveta",
             Inches(0.45), Inches(1.65), Inches(12), Inches(0.60),
             size_pt=28, bold=True, color=C_WHITE)
    add_text(slide, "Modelado SST-Captura y Diseño del Trigger  |  Perú",
             Inches(0.45), Inches(2.40), Inches(12), Inches(0.40),
             size_pt=14, color=C_MIDGRAY)
    add_text(slide, "Abril 2026   ·   Confidencial",
             Inches(0.45), Inches(2.90), Inches(12), Inches(0.35),
             size_pt=12, color=C_GRAY)


def section_slide(prs, number, title):
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_text(slide, number, Inches(0.55), Inches(1.10), Inches(4), Inches(1.20),
             size_pt=52, bold=True, color=C_GREEN)
    add_text(slide, title.upper(), Inches(0.55), Inches(2.40), Inches(12), Inches(0.60),
             size_pt=20, bold=True, color=C_WHITE)
    add_text(slide, "SUYANA", SLIDE_W - Inches(1.00), Inches(6.95),
             Inches(0.80), Inches(0.30), size_pt=7, bold=True,
             color=C_GREEN, align=PP_ALIGN.RIGHT)


def closing_slide(prs):
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_text(slide, "Gracias", Inches(0.45), Inches(1.60), Inches(12), Inches(0.90),
             size_pt=44, bold=True, color=C_WHITE)
    add_text(slide, "Seguro Paramétrico de Captura de Anchoveta Perú",
             Inches(0.45), Inches(2.60), Inches(12), Inches(0.40),
             size_pt=14, color=C_MIDGRAY)
    add_text(slide, "suyana.io   ·   Confidencial   ·   Abril 2026",
             Inches(0.45), Inches(3.10), Inches(12), Inches(0.35),
             size_pt=12, color=C_GRAY)
    add_text(slide, "SUYANA", SLIDE_W - Inches(1.00), Inches(6.95),
             Inches(0.80), Inches(0.30), size_pt=7, bold=True,
             color=C_GREEN, align=PP_ALIGN.RIGHT)


# ── Content slide builders ────────────────────────────────────────────────────

def s02_agenda(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Agenda")
    add_footer(slide, page_num)

    col_w = CONTENT_W / 2 - Inches(0.10)
    items_l = [("01", "Datos y cobertura"),
               ("02", "Contexto oceanográfico"),
               ("03", "Relación SST-Captura")]
    items_r = [("04", "Diseño del trigger"),
               ("05", "Probabilidad anual de excedencia"),
               ("06", "Análisis por empresa")]

    def agenda_col(x, items):
        y = ct + Inches(0.30)
        for num, desc in items:
            add_text(slide, num, x, y, Inches(0.55), Inches(0.38),
                     size_pt=16, bold=True, color=C_GREEN)
            add_text(slide, desc, x + Inches(0.60), y, col_w - Inches(0.60), Inches(0.38),
                     size_pt=16, color=C_BLACK)
            y += Inches(0.52)

    agenda_col(CONTENT_L, items_l)
    agenda_col(CONTENT_L + col_w + Inches(0.20), items_r)


def s04_fuentes(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Fuentes de datos",
                        "Tres fuentes integradas por interpolación espacio-temporal al nivel de cala")
    add_footer(slide, page_num)

    # Table headers
    headers = ["Fuente", "Variable principal", "Período", "Resolución"]
    rows = [
        ["IHMA (calas)", "Captura declarada (ton)", "2015-2024", "Por evento de pesca"],
        ["MODIS AQUA", "Anomalía SST diaria", "2002-2026", "~4 km diario"],
        ["HYCOM", "Temp. y salinidad subsuperficial", "2015-2024", "1/12° diario"],
    ]
    col_fracs = [0.18, 0.34, 0.18, 0.30]

    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h  = Inches(0.40)
    tbl_h  = n_rows * row_h
    tbl_w  = CONTENT_W

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        CONTENT_L, ct + Inches(0.15), tbl_w, tbl_h
    ).table

    for j, frac in enumerate(col_fracs):
        tbl.columns[j].width = int(tbl_w * frac)

    def set_cell(cell, text, bold=False, hdr=False):
        cell.text = text
        cell.fill.background()
        tf = cell.text_frame
        p  = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = bold
        p.font.color.rgb = C_WHITE if hdr else C_BLACK

    def cell_border(cell, pos, color_hex="141414", width_pt=1.0):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        ln   = etree.SubElement(tcPr, qn(f"a:{pos}"))
        ln.set("w", str(int(width_pt * 12700)))
        sf   = etree.SubElement(ln, qn("a:solidFill"))
        sc   = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", color_hex)

    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        set_cell(cell, hdr, bold=True, hdr=False)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_GREEN
        tf = cell.text_frame
        tf.paragraphs[0].font.color.rgb = C_WHITE
        cell_border(cell, "top", "43A047", 1.5)
        cell_border(cell, "bottom", "43A047", 1.0)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            set_cell(cell, val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            if i == len(rows) - 1:
                cell_border(cell, "bottom", "141414", 1.5)

    note_y = ct + Inches(0.15) + tbl_h + Inches(0.20)
    add_text(slide,
             "Variable final seleccionada: anomalía MODIS SST (T') — descartadas salinidad HYCOM "
             "(señal secundaria), clorofila (alta brecha por nubosidad) y temperatura "
             "subsuperficial HYCOM (colineal con SST MODIS).",
             CONTENT_L, note_y, CONTENT_W, Inches(0.60),
             size_pt=9, color=C_BLACK, wrap=True)


def s05_variables(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Variables incluidas y excluidas")
    add_footer(slide, page_num)

    col_w = CONTENT_W / 2 - Inches(0.10)

    def bullet_col(x, header, items, hdr_color=C_GREEN):
        y = ct + Inches(0.15)
        add_text(slide, header, x, y, col_w, Inches(0.30),
                 size_pt=12, bold=True, color=hdr_color)
        y += Inches(0.33)
        for item in items:
            add_text(slide, f"•  {item}", x + Inches(0.08), y,
                     col_w - Inches(0.08), Inches(0.26),
                     size_pt=11, color=C_BLACK, wrap=True)
            y += Inches(0.30)
        return y

    y = bullet_col(CONTENT_L,
                   "Incluida en el modelo",
                   ["Anomalía SST MODIS — señal más fuerte y robusta; cubre 2002-2026",
                    "Índice El Niño 1+2 — solo como contexto descriptivo (paso 09), no como regresor"])

    add_text(slide, "Decisión de temporalidad", CONTENT_L, y + Inches(0.10), col_w, Inches(0.30),
             size_pt=12, bold=True, color=C_GREEN)
    y += Inches(0.43)
    for item in ["Dos temporadas: T1 (abr-jul) y T2 (nov-dic)",
                 "Períodos fuera de temporada excluidos por regulación IMARPE"]:
        add_text(slide, f"•  {item}", CONTENT_L + Inches(0.08), y,
                 col_w - Inches(0.08), Inches(0.26), size_pt=11, color=C_BLACK, wrap=True)
        y += Inches(0.30)

    bullet_col(CONTENT_L + col_w + Inches(0.20),
               "Descartadas", [
                   "SST cruda (confunde efectos espaciales con anomalías)",
                   "Temp./salinidad HYCOM (colineal; menor cobertura)",
                   "Clorofila MODIS (alta patchiness; brechas por nubes)",
                   "Término cuadrático T'² (sin mejora en R² a nivel temporada)",
               ], hdr_color=C_BLACK)


def s06_density_season(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Densidad de calas por temporada (IHMA, 2015-2024)")
    add_footer(slide, page_num)
    cap_h = Inches(0.40)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step08_vessel_density_by_season.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Concentración de calas a lo largo del corredor peruano. T1 (abr-jul) y T2 (nov-dic).",
                top, ih, size_pt=8)


def s07_density_company(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Densidad de calas por empresa (IHMA, 2015-2024)")
    add_footer(slide, page_num)
    cap_h = Inches(0.40)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step08_vessel_density_by_company.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Distribución geográfica por empresa. Varias flotas concentran actividad "
                "en subzonas específicas del corredor.",
                top, ih, size_pt=8)


def s09_sst_t1(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Anomalía SST promedio T1 (MODIS AQUA, abr-jul, 2002-2024)")
    add_footer(slide, page_num)
    add_image_centered(slide, IMGDIR / "step17_sst_t1_maps.png", ct)


def s10_sst_t2(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Anomalía SST promedio T2 (MODIS AQUA, nov-dic, 2002-2024)")
    add_footer(slide, page_num)
    add_image_centered(slide, IMGDIR / "step17_sst_t2_maps.png", ct)


def s11_sst_ridge(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Distribución SST por banda de latitud (10 bandas, 2002-2026)")
    add_footer(slide, page_num)

    # Image: 82% width
    img_frac = 0.82
    avail_w  = CONTENT_W * img_frac
    avail_h  = FOOTER_RULE_Y - ct - Inches(0.10)
    iw, ih   = img_dims(IMGDIR / "step16_sst_ridge.png", avail_w, avail_h)
    slide.shapes.add_picture(str(IMGDIR / "step16_sst_ridge.png"),
                             CONTENT_L, ct + Emu(20000), iw, ih)

    # Text column: right side
    txt_x = CONTENT_L + avail_w + Inches(0.10)
    txt_w = CONTENT_W - avail_w - Inches(0.10)
    ty = ct + Inches(0.40)
    add_text(slide, "Resultados", txt_x, ty, txt_w, Inches(0.25),
             size_pt=9, bold=True, color=C_BLACK)
    ty += Inches(0.28)
    for item in ["285 meses MODIS AQUA",
                 "Bandas 8°-14°S: señal homogénea",
                 "El Niño 2015-16 y 2023 coherentes en toda la costa",
                 "7°-6°S más variable"]:
        add_text(slide, f"•  {item}", txt_x, ty, txt_w, Inches(0.34),
                 size_pt=8, color=C_BLACK, wrap=True)
        ty += Inches(0.36)
    ty += Inches(0.10)
    add_text(slide, "Implicancia: señal compartida - trigger único válido.",
             txt_x, ty, txt_w, Inches(0.50),
             size_pt=8, bold=True, color=C_BLACK, wrap=True)


def s12_sst_corr(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Coherencia espacial del SST: correlaciones entre bandas")
    add_footer(slide, page_num)

    # Image left: 55%
    avail_w = CONTENT_W * 0.55
    avail_h = FOOTER_RULE_Y - ct - Inches(0.10)
    iw, ih  = img_dims(IMGDIR / "step16_sst_correlations.png", avail_w, avail_h)
    slide.shapes.add_picture(str(IMGDIR / "step16_sst_correlations.png"),
                             CONTENT_L, ct + Emu(20000), iw, ih)

    # Text right: 41%
    txt_x = CONTENT_L + CONTENT_W * 0.57
    txt_w = CONTENT_W * 0.41
    ty = ct + Inches(0.15)
    add_text(slide, "Resultados principales", txt_x, ty, txt_w, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.32)
    for item in ["Mediana r = 0,81 (285 meses)",
                 "Mínimo r = 0,54 (7-6°S con 16-15°S)",
                 "Máximo r = 0,96 (10-9°S con 9-8°S)",
                 "80% de pares con r > 0,70"]:
        add_text(slide, f"•  {item}", txt_x, ty, txt_w, Inches(0.28),
                 size_pt=11, color=C_BLACK, wrap=True)
        ty += Inches(0.30)
    ty += Inches(0.10)
    add_text(slide,
             "Decisión: el calentamiento durante El Niño afecta simultáneamente toda la costa "
             "peruana. Un índice SST de área amplia captura la señal compartida sin necesidad "
             "de triggers diferenciados por subzona.",
             txt_x, ty, txt_w, Inches(0.90),
             size_pt=9, color=C_BLACK, wrap=True)


def s14_model_spec(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Especificación del modelo",
                        "Modelo M1: semi-log OLS captura total;  Modelo M2: semi-log OLS CPUE")
    add_footer(slide, page_num)

    col_w_l = CONTENT_W * 0.52
    col_w_r = CONTENT_W * 0.44

    ty = ct + Inches(0.10)
    # Left column: formulas as text
    add_text(slide, "Modelo M1 (captura)", CONTENT_L, ty, col_w_l, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.30)
    add_text(slide, "log(Cᵉ,s) = α + β · T'ᵉ,s + ε",
             CONTENT_L + Inches(0.30), ty, col_w_l, Inches(0.30),
             size_pt=13, italic=True, color=C_BLACK)
    ty += Inches(0.38)
    add_text(slide, "Modelo M2 (CPUE)", CONTENT_L, ty, col_w_l, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.30)
    add_text(slide, "log(Cᵉ,s / Eᵉ,s) = α + β · T'ᵉ,s + ε",
             CONTENT_L + Inches(0.30), ty, col_w_l, Inches(0.30),
             size_pt=13, italic=True, color=C_BLACK)
    ty += Inches(0.45)
    add_text(slide,
             "Símbolos: Cᵉ,s = captura (ton); Eᵉ,s = esfuerzo VMS (h); "
             "T'ᵉ,s = anomalía media SST (°C); α = intercepto; "
             "β = elasticidad; ε = residuo.",
             CONTENT_L, ty, col_w_l, Inches(0.70),
             size_pt=9, color=C_BLACK, wrap=True)

    # Right column
    rx = CONTENT_L + col_w_l + Inches(0.08)
    ry = ct + Inches(0.10)
    add_text(slide, "Decisiones clave", rx, ry, col_w_r, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ry += Inches(0.32)
    for item in ["empresa×temporada minimiza sesgo de atenuación",
                 "M1: 2015-2024;  M2: 2017-2022 (cobertura VMS)",
                 "Filtro IQR en calas individuales antes de agregar"]:
        add_text(slide, f"•  {item}", rx, ry, col_w_r, Inches(0.30),
                 size_pt=11, color=C_BLACK, wrap=True)
        ry += Inches(0.34)
    ry += Inches(0.08)
    add_text(slide,
             "Filtro IQR: calas con captura fuera de [Q1 - 1,5·IQR, Q3 + 1,5·IQR] "
             "excluidas antes de cualquier agregación.",
             rx, ry, col_w_r, Inches(0.55),
             size_pt=9, color=C_BLACK, wrap=True)


def s15_betas(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Beta por nivel de agregación (M1 y M2, región Centro)")
    add_footer(slide, page_num)
    cap_h = Inches(0.40)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step11_ols_betas.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Diario → temporada: beta se vuelve más negativo al reducirse el sesgo de "
                "atenuación. Empresa×temporada seleccionado.",
                top, ih, size_pt=8)


def s17_spatial(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Selección de la zona del trigger: comparación espacial")
    add_footer(slide, page_num)
    cap_h = Inches(0.50)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step15_spatial_comparison.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "OLS log(captura) ~ anomalía SST a nivel empresa×temporada.  "
                "Centro Norte (-11° a -7,1°S) tiene el mayor R² (0,261) y beta más negativo "
                "significativo (-0,816, p<0,001).  Norte: no significativo (p=0,058).  "
                "Centro Sur: sin relación (R²=0,002, p=0,65).",
                top, ih, size_pt=7)


def s18_payout_curves(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Curvas de pago: ramp lineal vs. step")
    add_footer(slide, page_num)
    cap_h = Inches(0.70)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step15_payout_curves.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Línea punteada gris = curva exponencial implícita por OLS (beta = -0,816, referencia).  "
                "Ramp lineal (recomendado): entrada en 0,5°C, máximo en 2,5°C.  "
                "Step: pago binario al p90 (0,96°C). El step sobreindemniza a anomalías bajas "
                "(100% al p90, OLS implica 62%); el ramp sigue más de cerca la pérdida real.",
                top, ih, size_pt=7)


def s19_formula(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Fórmula de pago: ramp lineal")
    add_footer(slide, page_num)

    col_w_l = CONTENT_W * 0.52
    col_w_r = CONTENT_W * 0.44

    ty = ct + Inches(0.10)
    add_text(slide, "Fracción de pago", CONTENT_L, ty, col_w_l, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.30)
    add_text(slide, "f = clip( (T' − T_ent) / (T_sal − T_ent),  0,  1 )",
             CONTENT_L + Inches(0.20), ty, col_w_l, Inches(0.30),
             size_pt=13, italic=True, color=C_BLACK)
    ty += Inches(0.40)
    add_text(slide, "Pago en toneladas", CONTENT_L, ty, col_w_l, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.30)
    add_text(slide, "P = B × f × c",
             CONTENT_L + Inches(0.20), ty, col_w_l, Inches(0.30),
             size_pt=13, italic=True, color=C_BLACK)
    ty += Inches(0.45)
    add_text(slide,
             "Símbolos: T' = anomalía media SST en la temporada (°C);  "
             "T_ent = 0,5°C = umbral de activación;  T_sal = 2,5°C = umbral de pago máximo;  "
             "f ∈ [0,1] = fracción del pago máximo;  "
             "B = captura de referencia de la empresa (ton);  "
             "c = cobertura contratada;  P = pago en ton equivalente.",
             CONTENT_L, ty, col_w_l, Inches(0.90),
             size_pt=9, color=C_BLACK, wrap=True)

    rx = CONTENT_L + col_w_l + Inches(0.08)
    ry = ct + Inches(0.10)
    add_text(slide, "Decisiones de diseño", rx, ry, col_w_r, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ry += Inches(0.32)
    for item in ["Beta OLS no entra en la fórmula; solo en calibración",
                 "B varía por empresa y temporada (T1 / T2)",
                 "Trigger SST = índice de área (toda la flota, Centro Norte): objetivo, no manipulable"]:
        add_text(slide, f"•  {item}", rx, ry, col_w_r, Inches(0.35),
                 size_pt=11, color=C_BLACK, wrap=True)
        ry += Inches(0.38)


def s20_examples(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Ejemplos de cálculo del pago",
                        "Tres regímenes históricos — fórmula aplicada con los mismos parámetros")
    add_footer(slide, page_num)

    col_w = (CONTENT_W - Inches(0.20)) / 3

    examples = [
        {
            "title": "Ej. 1 — T1-2019 (sin evento)",
            "subtitle": "SST normal, sin pago.",
            "items": [
                "T' = -0,3°C",
                "T' < T_ent = 0,5°C",
                "f = clip(...) = 0",
                "P = B × 0 × c",
                "Pago = 0",
            ],
            "note": "Sin pago; cuota regulatoria no está cubierta.",
        },
        {
            "title": "Ej. 2 — T2-2018 (falso trigger)",
            "subtitle": "SST leve; empresas sobre línea base.",
            "items": [
                "T' = +0,94°C",
                "f = 0,44 / 2,0 = 0,22",
                "P = 132.000 × 0,22 ≈ 29.000 ton",
                "22% pagado, sin pérdida real",
            ],
            "note": "Ramp limita sobreindemnización; step pagaría 100%.",
        },
        {
            "title": "Ej. 3 — T1-2023 (El Niño)",
            "subtitle": "Pérdida severa; alta cobertura.",
            "items": [
                "T' = +2,01°C",
                "f = 1,51 / 2,0 = 0,755",
                "P = 132.000 × 0,755 ≈ 100.000 ton",
                "Pérdida real: ~124.000 ton",
                "Cobertura 80%",
            ],
            "note": "Descubierto (20%) porque T' < T_sal = 2,5.",
        },
    ]

    for i, ex in enumerate(examples):
        x = CONTENT_L + i * (col_w + Inches(0.10))
        ty = ct + Inches(0.05)
        add_text(slide, ex["title"], x, ty, col_w, Inches(0.28),
                 size_pt=11, bold=True, color=C_BLACK)
        ty += Inches(0.30)
        add_text(slide, ex["subtitle"], x, ty, col_w, Inches(0.22),
                 size_pt=9, italic=True, color=C_GRAY)
        ty += Inches(0.25)
        for j, item in enumerate(ex["items"], 1):
            is_last = (j == len(ex["items"]))
            add_text(slide, f"{j}. {item}", x + Inches(0.05), ty,
                     col_w - Inches(0.05), Inches(0.26),
                     size_pt=9, bold=is_last, color=C_GREEN if is_last else C_BLACK,
                     wrap=True)
            ty += Inches(0.28)
        ty += Inches(0.08)
        add_text(slide, ex["note"], x, ty, col_w, Inches(0.45),
                 size_pt=9, italic=True, color=C_GRAY, wrap=True)


def s22_sst_timeseries(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Serie temporal de anomalía SST por temporada")
    add_footer(slide, page_num)
    cap_h = Inches(0.45)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step13_sst_timeseries.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Anomalía SST media por temporada (Centro Norte, 2002-2026, MODIS AQUA). "
                "Se marcan los percentiles p90/p95/p99 calibrados sobre el período completo. "
                "El Niño 2015-2016 y 2023 superan el p95.",
                top, ih, size_pt=8)


def s23_aep(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Curva de probabilidad anual de excedencia")
    add_footer(slide, page_num)

    # Image left: 58%
    avail_w = CONTENT_W * 0.58
    avail_h = FOOTER_RULE_Y - ct - Inches(0.10)
    iw, ih  = img_dims(IMGDIR / "step13_bootstrap_aep.png", avail_w, avail_h)
    slide.shapes.add_picture(str(IMGDIR / "step13_bootstrap_aep.png"),
                             CONTENT_L, ct + Emu(20000), iw, ih)

    # Text right: 38%
    tx = CONTENT_L + CONTENT_W * 0.60
    tw = CONTENT_W * 0.38
    ty = ct + Inches(0.10)
    add_text(slide, "Metodología", tx, ty, tw, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.32)
    for item in ["Ajuste Normal a anomalías SST estacionales",
                 "4.000 simulaciones bootstrap; IC 90%"]:
        add_text(slide, f"•  {item}", tx, ty, tw, Inches(0.28),
                 size_pt=11, color=C_BLACK, wrap=True)
        ty += Inches(0.30)
    ty += Inches(0.10)
    add_text(slide, "Percentiles del trigger", tx, ty, tw, Inches(0.28),
             size_pt=12, bold=True, color=C_BLACK)
    ty += Inches(0.32)
    for item in ["p90 = 0,96°C (~10 años de retorno)",
                 "p95 = 1,38°C (~20 años)",
                 "p99 = 2,75°C (~100 años)"]:
        add_text(slide, f"•  {item}", tx, ty, tw, Inches(0.28),
                 size_pt=11, color=C_BLACK, wrap=True)
        ty += Inches(0.30)
    ty += Inches(0.10)
    add_text(slide,
             "T1 ≈ 0,97 M ton;  T2 ≈ 1,27 M ton;  total ≈ 2,24 M ton (línea base, Centro).",
             tx, ty, tw, Inches(0.40),
             size_pt=9, color=C_BLACK, wrap=True)


def s25_companies(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "COPEINCA y EXALMAR: análisis contrafactual")
    add_footer(slide, page_num)
    cap_h = Inches(0.70)
    top, _, ih = add_image_centered(slide,
                                    IMGDIR / "step18_client_copeinca_exalmar_nosc.png",
                                    ct, cap_h_emu=cap_h)
    add_caption(slide,
                "Barras azules: captura real; rojo: temporadas bajo línea base; amarillo: pago simulado.  "
                "COPEINCA LB T1=132 kt, T2=118 kt.    EXALMAR LB T1=52 kt, T2=49 kt.  "
                "Panel inferior: anomalía SST estacional con umbrales de entrada y salida del trigger.",
                top, ih, size_pt=7)


def s26_base_risk(prs, page_num):
    slide = blank_slide(prs)
    add_header(slide)
    ct = add_frametitle(slide, "Riesgo base: observaciones clave",
                        "Patrones donde el trigger SST y la pérdida real divergen")
    add_footer(slide, page_num)

    col_w = CONTENT_W / 2 - Inches(0.10)

    ty = ct + Inches(0.10)
    add_text(slide, "T2-2018: falso trigger", CONTENT_L, ty, col_w, Inches(0.28),
             size_pt=11, bold=True, color=C_GREEN)
    ty += Inches(0.30)
    for item in ["SST = +0,94°C, f = 22%; sin pérdida real",
                 "COPEINCA +38%, EXALMAR +71% sobre línea base"]:
        add_text(slide, f"•  {item}", CONTENT_L, ty, col_w, Inches(0.26),
                 size_pt=11, color=C_BLACK, wrap=True)
        ty += Inches(0.28)
    ty += Inches(0.12)
    add_text(slide, "T1-2023: El Niño severo", CONTENT_L, ty, col_w, Inches(0.28),
             size_pt=11, bold=True, color=C_GREEN)
    ty += Inches(0.30)
    for item in ["SST = +2,01°C, f = 75,5%",
                 "COPEINCA 6%, EXALMAR 10% de línea base",
                 "Cobertura ~80%; descubierto 20%"]:
        add_text(slide, f"•  {item}", CONTENT_L, ty, col_w, Inches(0.26),
                 size_pt=11, color=C_BLACK, wrap=True)
        ty += Inches(0.28)

    rx = CONTENT_L + col_w + Inches(0.20)
    ry = ct + Inches(0.10)
    add_text(slide, "Riesgo base estructural", rx, ry, col_w, Inches(0.28),
             size_pt=11, bold=True, color=C_BLACK)
    ry += Inches(0.30)
    for item in ["Pérdidas en años fríos (cuotas IMARPE) no cubiertas",
                 "R² empresa 0,26-0,40: factores no-SST explican parte de la varianza",
                 "Ramp reduce sobreindemnización vs. step"]:
        add_text(slide, f"•  {item}", rx, ry, col_w, Inches(0.26),
                 size_pt=11, color=C_BLACK, wrap=True)
        ry += Inches(0.30)
    ry += Inches(0.15)
    add_text(slide,
             "El producto cubre bien El Niño severo. Riesgo base moderado "
             "(T' en [0,5, 1,5]°C) es inherente a la estructura paramétrica.",
             rx, ry, col_w, Inches(0.55),
             size_pt=9, color=C_BLACK, wrap=True)


# ── Main builder ─────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: Title
    title_slide(prs, 1)

    # Slide 2: Agenda
    s02_agenda(prs, 2)

    # Slide 3: §01 section
    section_slide(prs, "01", "Datos y Cobertura")

    # Slide 4: Fuentes de datos
    s04_fuentes(prs, 4)

    # Slide 5: Variables incluidas/excluidas
    s05_variables(prs, 5)

    # Slide 6: Densidad por temporada
    s06_density_season(prs, 6)

    # Slide 7: Densidad por empresa
    s07_density_company(prs, 7)

    # Slide 8: §02 section
    section_slide(prs, "02", "Contexto Oceanográfico")

    # Slide 9: SST T1 maps
    s09_sst_t1(prs, 9)

    # Slide 10: SST T2 maps
    s10_sst_t2(prs, 10)

    # Slide 11: SST ridge
    s11_sst_ridge(prs, 11)

    # Slide 12: SST correlations
    s12_sst_corr(prs, 12)

    # Slide 13: §03 section
    section_slide(prs, "03", "Relación SST-Captura")

    # Slide 14: Model specification
    s14_model_spec(prs, 14)

    # Slide 15: Betas
    s15_betas(prs, 15)

    # Slide 16: §04 section
    section_slide(prs, "04", "Diseño del Trigger")

    # Slide 17: Spatial comparison
    s17_spatial(prs, 17)

    # Slide 18: Payout curves
    s18_payout_curves(prs, 18)

    # Slide 19: Payout formula
    s19_formula(prs, 19)

    # Slide 20: Calculation examples
    s20_examples(prs, 20)

    # Slide 21: §05 section
    section_slide(prs, "05", "Probabilidad Anual de Excedencia")

    # Slide 22: SST time series
    s22_sst_timeseries(prs, 22)

    # Slide 23: AEP curve
    s23_aep(prs, 23)

    # Slide 24: §06 section
    section_slide(prs, "06", "Análisis por Empresa")

    # Slide 25: Company analysis
    s25_companies(prs, 25)

    # Slide 26: Base risk
    s26_base_risk(prs, 26)

    # Slide 27: Closing
    closing_slide(prs)

    prs.save(OUT)
    print(f"Saved {OUT}  ({TOTAL} slides)")


if __name__ == "__main__":
    build()
